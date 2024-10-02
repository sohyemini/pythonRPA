from pptx import Presentation

prs = Presentation() # make a presentation
title_slide_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save(r"./files/test.pptx")